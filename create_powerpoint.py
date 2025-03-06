from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()

# Function to add a slide with a title and content
def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]  # Layout for title and content
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_placeholder = slide.shapes.title
    title_placeholder.text = title

    # Content
    content_placeholder = slide.shapes.placeholders[1]
    text_frame = content_placeholder.text_frame
    text_frame.text = content

    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = RGBColor(0, 0, 0)
        paragraph.alignment = PP_ALIGN.CENTER

# Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.shapes.placeholders[1]

title.text = "Chocolate Pooping Elf Storyboard"
subtitle.text = "Visual Storyboard for Concept Development"

# Slide 1: Idea 1 - Introduction of the Elf
add_slide(
    "Idea 1: Introduction of the Elf",
    "The elf is cheerful and lives in a magical forest. He has the unique ability to poop chocolate! "
    "Show the elf in his natural environment, surrounded by trees and magic sparkles."
)

# Slide 2: Idea 2 - Chocolate Factory Scene
add_slide(
    "Idea 2: Chocolate Factory Scene",
    "The elf visits a factory, where his chocolate is refined and packaged into gourmet treats. "
    "Visualize the factory scene, with conveyor belts and excited workers packaging the chocolate."
)

# Slide 3: Idea 3 - Marketing and Branding
add_slide(
    "Idea 3: Marketing and Branding",
    "Create a fun and whimsical brand around the elf's chocolate. The logo might show a playful elf "
    "with a mischievous grin, perhaps holding a chocolate bar. Include slogans like 'Magically Delicious!'"
)

# Slide 4: Idea 4 - Final Delivery Scene
add_slide(
    "Idea 4: Final Delivery Scene",
    "The elf personally delivers his chocolate treats to kids and families, bringing joy and laughter. "
    "Show happy families unwrapping the chocolate and enjoying it."
)

# Slide 5: Ending and Next Steps
add_slide(
    "Ending: Next Steps",
    "Summarize the concept and outline the next steps for production. These may include animation, voiceover, "
    "and packaging design. End with a call to action for further development."
)

# Save the presentation
prs.save("Chocolate_Pooping_Elf_Storyboard.pptx")
print("Presentation saved as 'Chocolate_Pooping_Elf_Storyboard.pptx'.")